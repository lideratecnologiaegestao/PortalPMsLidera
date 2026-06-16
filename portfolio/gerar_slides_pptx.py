# -*- coding: utf-8 -*-
"""
Gera a APRESENTAÇÃO do "Portal de Prefeitura" em .pptx (16:9), pronta para
importar no Google Apresentações (Google Slides).
Dependências: python-pptx (e os PNGs em portfolio/assets/, gerados por
gerar_portfolio_docx.py — rode aquele script antes, se os assets não existirem).
Uso: python gerar_slides_pptx.py
Saída: portfolio/Apresentacao-Portal-Prefeitura.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# paleta gov.br
AZUL      = RGBColor(0x13, 0x51, 0xB4)
AZUL_ESC  = RGBColor(0x0C, 0x32, 0x6F)
VERDE     = RGBColor(0x16, 0x88, 0x21)
AMARELO   = RGBColor(0xFF, 0xCD, 0x07)
BRANCO    = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_BG  = RGBColor(0xF2, 0xF5, 0xFA)
CINZA_TX  = RGBColor(0x33, 0x39, 0x42)
CINZA_MD  = RGBColor(0x55, 0x60, 0x70)

BASE   = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(BASE, "assets")
LOGOS  = os.path.join(BASE, "Logos")
LOGO_PNG     = os.path.join(ASSETS, "logo_lidera.png")        # colorido (fundo claro)
LOGO_DARK    = os.path.join(ASSETS, "logo_lidera_dark.png")   # branco sobre #0C326F (slides escuros)

EMAIL = "comercial@lideratecnologia.com.br"
SITE  = "www.lideratecnologia.com.br"
DEMO  = "prefeitura.lidera.app.br"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def textbox(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    """runs: lista de parágrafos; cada parágrafo é lista de (texto, size, color, bold, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        for (txt, size, color, bold, italic) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic
            r.font.name = "Calibri"
    return tb

def pic_contain(s, path, x, y, w, h, align="center", valign="middle"):
    """Insere imagem mantendo proporção, contida na caixa (x,y,w,h)."""
    from PIL import Image
    iw, ih = Image.open(path).size
    bw, bh = w, h
    scale = min(bw / iw, bh / ih)
    nw, nh = Emu(int(iw * scale)), Emu(int(ih * scale))
    px = x + (Emu(0) if align == "left" else (bw - nw) // 2 if align == "center" else bw - nw)
    py = y + ((bh - nh) // 2 if valign == "middle" else Emu(0) if valign == "top" else bh - nh)
    return s.shapes.add_picture(path, px, py, nw, nh)

def header(s, kicker, title):
    """Cabeçalho padrão de slide de conteúdo (fundo branco)."""
    rect(s, 0, 0, SW, SH, BRANCO)
    rect(s, 0, 0, Inches(0.18), SH, AZUL)            # faixa lateral
    rect(s, Inches(0.6), Inches(0.55), Inches(0.9), Inches(0.06), AMARELO)
    textbox(s, Inches(0.6), Inches(0.62), Inches(11), Inches(0.4),
            [[(kicker.upper(), 12, AZUL, True, False)]])
    textbox(s, Inches(0.58), Inches(0.95), Inches(12), Inches(0.9),
            [[(title, 30, AZUL_ESC, True, False)]])

def footer(s, dark=False):
    col = RGBColor(0xA9,0xB6,0xCC) if dark else RGBColor(0x9A,0xA6,0xB8)
    textbox(s, Inches(0.6), SH - Inches(0.45), Inches(12.1), Inches(0.3),
            [[("Lidera Tecnologia e Gestão  ·  Portal de Prefeitura  ·  " + SITE, 9, col, False, False)]],
            align=PP_ALIGN.LEFT)

def bullets(s, items, x=Inches(0.7), y=Inches(2.0), w=Inches(11.9), h=Inches(4.8), size=16, gap=8):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        r = p.add_run(); r.text = "▸  "; r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = AZUL; r.font.name = "Calibri"
        r2 = p.add_run(); r2.text = it; r2.font.size = Pt(size)
        r2.font.color.rgb = CINZA_TX; r2.font.name = "Calibri"
    return tb

# =====================================================================
# 1 — CAPA
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, AZUL_ESC)
rect(s, 0, 0, SW, Inches(0.16), AMARELO)
rect(s, 0, SH - Inches(0.16), SW, Inches(0.16), VERDE)
pic_contain(s, LOGO_DARK, Inches(5.52), Inches(0.6), Inches(2.3), Inches(2.15))
textbox(s, Inches(1), Inches(3.0), Inches(11.3), Inches(1.6),
        [[("Portal de Prefeitura", 46, BRANCO, True, False)]], align=PP_ALIGN.CENTER)
textbox(s, Inches(1.4), Inches(4.3), Inches(10.5), Inches(1.0),
        [[("A plataforma digital completa, acessível e em conformidade com a lei para o seu município.",
           18, RGBColor(0xDD,0xE6,0xF6), False, False)]], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(5.7), Inches(11.3), Inches(1.0),
        [[("Selo Diamante (PNTP) · Multi-tenant · Gov.br · LGPD", 14, AMARELO, True, False)],
         [(f"{EMAIL}   ·   {SITE}   ·   Demo: {DEMO}", 13, RGBColor(0xCF,0xDA,0xEE), False, False)]],
        align=PP_ALIGN.CENTER)

# =====================================================================
# 2 — AGENDA
# =====================================================================
s = slide(); header(s, "Agenda", "O que você vai ver")
col1 = ["O desafio das prefeituras hoje", "A plataforma e a arquitetura",
        "Modelo SaaS multi-tenant", "Catálogo de 20+ módulos", "Conformidade legal nativa"]
col2 = ["Transparência: 100% PNTP (Diamante)", "Segurança, privacidade e LGPD",
        "App do Cidadão e Inteligência Artificial", "Caso real e implantação", "Próximos passos"]
bullets(s, col1, x=Inches(0.7),  y=Inches(2.1), w=Inches(6), size=17, gap=12)
bullets(s, col2, x=Inches(6.9),  y=Inches(2.1), w=Inches(6), size=17, gap=12)
footer(s)

# =====================================================================
# 3 — O DESAFIO
# =====================================================================
s = slide(); header(s, "O desafio", "Muita exigência legal, pouco orçamento de TI")
bullets(s, [
 "Transparência e PNTP: difícil atingir e manter o selo da Atricon com dados dispersos.",
 "Prazos legais: e-SIC (LAI) e Ouvidoria (Lei 13.460) exigem controle que planilha não dá.",
 "Acessibilidade: a lei exige WCAG/eMAG e VLibras — a maioria dos sites não cumpre.",
 "Custo e dependência: cada sistema é um contrato e um fornecedor diferente.",
 "LGPD: dados do cidadão tratados sem base legal, sem direitos do titular, sem registro de incidentes.",
], size=17, gap=14)
footer(s)

# =====================================================================
# 4 — A SOLUÇÃO
# =====================================================================
s = slide(); header(s, "A solução", "Uma plataforma, completa e em conformidade")
textbox(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.9),
        [[("Site institucional, transparência, ouvidoria/e-SIC, diário oficial, serviços ao cidadão, "
           "app mobile e inteligência artificial — em um só produto, com a legislação embutida.",
           17, CINZA_TX, False, False)]])
cards = [("Tudo integrado", "20+ módulos no mesmo login e nos mesmos dados", AZUL),
         ("Conforme a lei", "LAI · Lei 13.460 · LC 131/LRF · LGPD · WCAG", VERDE),
         ("100% PNTP", "Selo Diamante em transparência pública", AMARELO),
         ("Multi-tenant", "uma infra serve N prefeituras, com identidade própria", AZUL_ESC)]
cx = Inches(0.7); cw = Inches(2.95); gap = Inches(0.13)
for i,(t,d,c) in enumerate(cards):
    x = cx + (cw + gap) * i
    rect(s, x, Inches(3.1), cw, Inches(2.6), CINZA_BG)
    rect(s, x, Inches(3.1), cw, Inches(0.13), c)
    tc = AZUL_ESC if c == AMARELO else c
    textbox(s, x+Inches(0.18), Inches(3.4), cw-Inches(0.36), Inches(0.7),
            [[(t, 17, tc, True, False)]])
    textbox(s, x+Inches(0.18), Inches(4.1), cw-Inches(0.36), Inches(1.5),
            [[(d, 13, CINZA_TX, False, False)]])
footer(s)

# =====================================================================
# 5 — ARQUITETURA (imagem)
# =====================================================================
s = slide(); header(s, "Plataforma", "Arquitetura segura por projeto")
pic_contain(s, os.path.join(ASSETS, "arquitetura.png"), Inches(0.7), Inches(1.85),
            Inches(8.2), Inches(5.0), align="center")
bullets(s, ["Frontend e app falam só com a API (gateway único).",
            "Dados isolados por RLS — prefeitura nunca vê dados de outra.",
            "Duas camadas: RBAC (o que pode fazer) + RLS (o que pode ver).",
            "Tecnologia moderna, aberta e auditável."],
        x=Inches(9.1), y=Inches(2.2), w=Inches(3.9), size=13, gap=12)
footer(s)

# =====================================================================
# 6 — MULTI-TENANT (imagem)
# =====================================================================
s = slide(); header(s, "Modelo SaaS", "Uma plataforma, muitas prefeituras")
pic_contain(s, os.path.join(ASSETS, "multitenant.png"), Inches(0.8), Inches(1.95),
            Inches(11.7), Inches(3.3), align="center")
bullets(s, ["Custo compartilhado entre as prefeituras atendidas.",
            "Atualizou uma vez → todas recebem (correções legais inclusas).",
            "Identidade própria (white-label): cores, logo, domínio e conteúdo."],
        x=Inches(0.9), y=Inches(5.4), w=Inches(11.5), size=15, gap=8)
footer(s)

# =====================================================================
# 7 — MÓDULOS
# =====================================================================
s = slide(); header(s, "Produto", "20+ módulos integrados, em 6 áreas")
grupos = [("Cidadão & Participação", AZUL, "Ouvidoria/e-SIC · Atendimento omnichannel · Enquetes · Formulários · App do Cidadão"),
          ("Transparência & Contas", VERDE, "Portal da Transparência · APLIC/TCE-MT · PNTP Diamante · Diário Oficial"),
          ("Conteúdo & Comunicação", RGBColor(0xB3,0x53,0x1D), "CMS drag-drop · Notícias · Secretarias · Galeria · Documentos · Carta de Serviços"),
          ("Gestão & Administração", RGBColor(0x6f,0x42,0xc1), "Gerenciador multi-tenant · Usuários/Grupos/Sessões · Chat interno · Configurações"),
          ("Inteligência Artificial", RGBColor(0x15,0x5F,0x8A), "Chatbot RAG · Triagem de manifestações · Busca com OCR · IA fiscal"),
          ("Conformidade & Segurança", RGBColor(0xE5,0x22,0x07), "LGPD self-service · RLS · RBAC · WCAG/VLibras · Login gov.br")]
gx = Inches(0.65); gw = Inches(3.95); gh = Inches(2.15); gxg = Inches(0.1); gyg = Inches(0.16)
for i,(t,c,d) in enumerate(grupos):
    col = i % 3; row = i // 3
    x = gx + (gw + gxg) * col
    y = Inches(2.0) + (gh + gyg) * row
    rect(s, x, y, gw, gh, CINZA_BG)
    rect(s, x, y, gw, Inches(0.5), c)
    textbox(s, x+Inches(0.18), y+Inches(0.05), gw-Inches(0.36), Inches(0.45),
            [[(t, 14, BRANCO, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, x+Inches(0.18), y+Inches(0.62), gw-Inches(0.36), Inches(1.45),
            [[(d, 12, CINZA_TX, False, False)]])
footer(s)

# =====================================================================
# 8 — CONFORMIDADE (imagem)
# =====================================================================
s = slide(); header(s, "Conformidade", "A lei está no núcleo do produto")
pic_contain(s, os.path.join(ASSETS, "conformidade.png"), Inches(0.7), Inches(1.9),
            Inches(7.7), Inches(4.9), align="center")
bullets(s, ["Prazos legais já configurados: e-SIC 20+10, Ouvidoria 30+30.",
            "O sistema alerta e vence os prazos automaticamente.",
            "Acessibilidade obrigatória: tema reprovado no contraste não salva.",
            "Login do cidadão via gov.br (Login Único)."],
        x=Inches(8.7), y=Inches(2.2), w=Inches(4.3), size=14, gap=12)
footer(s)

# =====================================================================
# 9 — PNTP 100% (destaque)
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, AZUL_ESC)
rect(s, 0, 0, SW, Inches(0.14), AMARELO)
textbox(s, Inches(0.7), Inches(0.55), Inches(12), Inches(0.4),
        [[("TRANSPARÊNCIA · PNTP / ATRICON", 13, AMARELO, True, False)]])
textbox(s, Inches(0.68), Inches(0.95), Inches(12), Inches(0.9),
        [[("100% dos critérios — Selo Diamante", 32, BRANCO, True, False)]])
# painel branco com o gráfico
rect(s, Inches(0.7), Inches(2.0), Inches(11.93), Inches(4.6), BRANCO)
pic_contain(s, os.path.join(ASSETS, "pntp.png"), Inches(0.9), Inches(2.2),
            Inches(11.5), Inches(3.5), align="center")
textbox(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.6),
        [[("Sua prefeitura entra no ar já na nota máxima da transparência pública — "
           "Receita, Despesa, RH, Licitações e Contratos em 100%.", 14, CINZA_TX, False, False)]],
        align=PP_ALIGN.CENTER)

# =====================================================================
# 10 — SEGURANÇA
# =====================================================================
s = slide(); header(s, "Segurança & Privacidade", "Dados protegidos e isolados")
bullets(s, [
 "Isolamento por RLS: cada consulta é restrita ao município.",
 "RBAC: papéis (administrador, gestor, ouvidor, servidor, cidadão) controlam cada ação.",
 "LGPD por projeto: base legal por finalidade e logs de acesso a dados pessoais.",
 "Auditoria de toda ação sensível e de falhas de processamento.",
 "Borda protegida: nada exposto direto à internet — proxy + WAF + TLS.",
 "Login gov.br: identidade forte sem o município guardar senha.",
], size=16, gap=11)
footer(s)

# =====================================================================
# 11 — APP + IA
# =====================================================================
s = slide(); header(s, "Cidadão & IA", "App do Cidadão + Inteligência Artificial")
rect(s, Inches(0.7), Inches(2.0), Inches(5.85), Inches(4.5), CINZA_BG)
rect(s, Inches(0.7), Inches(2.0), Inches(5.85), Inches(0.5), AZUL)
textbox(s, Inches(0.9), Inches(2.05), Inches(5.5), Inches(0.45),
        [[("📱  App do Cidadão (mobile)", 16, BRANCO, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
bullets(s, ["Chamados georreferenciados: buracos, terrenos, animais.",
            "Foto + GPS, detecção de duplicados por proximidade.",
            "Acompanhamento e notificações (push).",
            "Login gov.br · tema da prefeitura (white-label)."],
        x=Inches(0.95), y=Inches(2.7), w=Inches(5.4), size=13.5, gap=10)
rect(s, Inches(6.75), Inches(2.0), Inches(5.85), Inches(4.5), CINZA_BG)
rect(s, Inches(6.75), Inches(2.0), Inches(5.85), Inches(0.5), VERDE)
textbox(s, Inches(6.95), Inches(2.05), Inches(5.5), Inches(0.45),
        [[("🤖  Inteligência Artificial", 16, BRANCO, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
bullets(s, ["Chatbot que responde da base oficial da prefeitura (RAG).",
            "Triagem automática de manifestações, com revisão humana.",
            "Busca unificada com OCR de PDFs digitalizados.",
            "IA fiscal sobre os dados da transparência."],
        x=Inches(7.0), y=Inches(2.7), w=Inches(5.4), size=13.5, gap=10)
footer(s)

# =====================================================================
# 12 — CASO REAL
# =====================================================================
s = slide(); header(s, "Prova de entrega", "Caso real — Barão de Melgaço (MT)")
textbox(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.8),
        [[("Site legado (Joomla) migrado integralmente para a plataforma, preservando histórico e SEO:",
           16, CINZA_TX, False, False)]])
nums = [("455", "notícias migradas"), ("1.180", "documentos migrados"),
        ("13", "secretarias publicadas"), ("100%", "critérios PNTP (Diamante)")]
nx = Inches(0.7); nw = Inches(2.95); ng = Inches(0.13)
for i,(n,d) in enumerate(nums):
    x = nx + (nw + ng) * i
    rect(s, x, Inches(2.9), nw, Inches(2.0), CINZA_BG)
    textbox(s, x, Inches(3.1), nw, Inches(1.0), [[(n, 40, AZUL, True, False)]], align=PP_ALIGN.CENTER)
    textbox(s, x, Inches(4.15), nw, Inches(0.7), [[(d, 13, CINZA_TX, False, False)]], align=PP_ALIGN.CENTER)
textbox(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.0),
        [[("Também migrados o institucional, a galeria e a Carta de Serviços — com redirecionamentos "
           "que preservam os links indexados pelo Google.", 14, CINZA_MD, False, True)]])
footer(s)

# =====================================================================
# 13 — IMPLANTAÇÃO + ROADMAP
# =====================================================================
s = slide(); header(s, "Implantação", "Roda onde a prefeitura precisar")
bullets(s, ["Nuvem gerenciada (Google Cloud / AWS) — provisionada por Terraform.",
            "On-premise (Windows / Linux / Docker) — no servidor da prefeitura.",
            "Migração do site atual sem perder histórico nem SEO.",
            "Treinamento das equipes e suporte continuado."],
        x=Inches(0.7), y=Inches(2.0), w=Inches(12), size=16, gap=10)
pic_contain(s, os.path.join(ASSETS, "roadmap.png"), Inches(0.8), Inches(4.7),
            Inches(11.7), Inches(2.2), align="center")
footer(s)

# =====================================================================
# 14 — ENCERRAMENTO / CTA
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, AZUL_ESC)
rect(s, 0, 0, SW, Inches(0.16), AMARELO)
rect(s, 0, SH - Inches(0.16), SW, Inches(0.16), VERDE)
pic_contain(s, LOGO_DARK, Inches(5.92), Inches(0.55), Inches(1.5), Inches(1.6))
textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.0),
        [[("Vamos modernizar a sua prefeitura?", 34, BRANCO, True, False)]], align=PP_ALIGN.CENTER)
textbox(s, Inches(1.4), Inches(3.7), Inches(10.5), Inches(0.9),
        [[("Agende uma demonstração gratuita e veja o portal funcionando com os dados do seu município.",
           17, RGBColor(0xDD,0xE6,0xF6), False, False)]], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.9), Inches(11.3), Inches(1.3),
        [[(f"✉  {EMAIL}", 20, AMARELO, True, False)],
         [(f"🌐  {SITE}      🔗  Demonstração ao vivo: {DEMO}", 15, RGBColor(0xCF,0xDA,0xEE), False, False)]],
        align=PP_ALIGN.CENTER)
textbox(s, Inches(1), SH - Inches(0.7), Inches(11.3), Inches(0.4),
        [[("Lidera Tecnologia e Gestão  ·  Investimento sob consulta, conforme o porte do município.",
           11, RGBColor(0xA9,0xB6,0xCC), False, True)]], align=PP_ALIGN.CENTER)

out = os.path.join(BASE, "Apresentacao-Portal-Prefeitura.pptx")
prs.save(out)
print(f"OK -> {out}  ({len(prs.slides._sldIdLst)} slides)")
